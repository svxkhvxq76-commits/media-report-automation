#!/usr/bin/env python3
"""Pipeline robusto de análise de investimento por canal com exportação para PowerPoint."""

from __future__ import annotations

import argparse
import csv
import json
import logging
from collections import defaultdict
from dataclasses import asdict, dataclass
from datetime import datetime
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

REQUIRED_COLUMNS = ["Canal", "Investimento"]


@dataclass
class KPIBundle:
    total_investimento: float
    canais_ativos: int
    investimento_medio_por_canal: float
    top_canal: str
    top_canal_investimento: float
    top_canal_share: float
    hhi_concentracao: float


@dataclass
class RankingRow:
    canal: str
    investimento: float
    share: float
    share_percentual: float
    investimento_acumulado: float
    share_acumulado: float


@dataclass
class ReportArtifacts:
    base_processada_csv: Path
    ranking_canais_csv: Path
    kpis_json: Path
    kpis_csv: Path
    grafico_pizza_png: Path
    grafico_barras_png: Path
    apresentacao_pptx: Path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Lê uma planilha (.xlsx) ou CSV com Canal e Investimento, calcula KPIs, "
            "gera visualizações e exporta um relatório em PowerPoint."
        )
    )
    parser.add_argument("--input", required=True, help="Caminho do arquivo de entrada (.xlsx ou .csv).")
    parser.add_argument("--sheet", default="Resumo", help="Nome da aba para .xlsx (padrão: Resumo).")
    parser.add_argument("--output-dir", default="output", help="Diretório para arquivos gerados.")
    parser.add_argument("--title", default="Share de Investimento por Canal", help="Título principal do relatório.")
    parser.add_argument("--campaign", default="Campanha não informada", help="Nome da campanha para contextualização.")
    parser.add_argument("--period", default="Período não informado", help="Período da campanha (ex.: Jan/2026).")
    parser.add_argument("--log-level", default="INFO", choices=["DEBUG", "INFO", "WARNING", "ERROR"])
    return parser.parse_args()


def setup_logging(level: str) -> None:
    logging.basicConfig(level=getattr(logging, level), format="%(asctime)s | %(levelname)s | %(message)s")


def _load_rows_from_csv(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as fp:
        reader = csv.DictReader(fp)
        return [dict(row) for row in reader]


def _load_rows_from_xlsx(path: Path, sheet_name: str) -> list[dict[str, str]]:
    try:
        from openpyxl import load_workbook
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError("Para ler .xlsx instale openpyxl: pip install openpyxl") from exc

    workbook = load_workbook(path, data_only=True)
    if sheet_name not in workbook.sheetnames:
        raise ValueError(f"Aba '{sheet_name}' não encontrada. Abas disponíveis: {workbook.sheetnames}")

    sheet = workbook[sheet_name]
    rows = list(sheet.iter_rows(values_only=True))
    if not rows:
        return []

    headers = [str(h).strip() if h is not None else "" for h in rows[0]]
    out: list[dict[str, str]] = []
    for values in rows[1:]:
        record = {}
        for idx, header in enumerate(headers):
            if not header:
                continue
            record[header] = "" if idx >= len(values) or values[idx] is None else str(values[idx]).strip()
        out.append(record)
    return out


def _validate_columns(rows: list[dict[str, str]]) -> None:
    if not rows:
        raise ValueError("Arquivo sem linhas de dados.")
    header_set = set(rows[0].keys())
    missing = [column for column in REQUIRED_COLUMNS if column not in header_set]
    if missing:
        raise ValueError(f"Colunas obrigatórias ausentes: {missing}")


def load_and_validate_data(input_path: Path, sheet_name: str) -> list[dict[str, float | str]]:
    logging.info("Lendo arquivo: %s", input_path)

    suffix = input_path.suffix.lower()
    if suffix == ".csv":
        raw_rows = _load_rows_from_csv(input_path)
    elif suffix == ".xlsx":
        raw_rows = _load_rows_from_xlsx(input_path, sheet_name)
    else:
        raise ValueError("Formato não suportado. Use .csv ou .xlsx")

    _validate_columns(raw_rows)

    cleaned: list[dict[str, float | str]] = []
    for row in raw_rows:
        canal = str(row.get("Canal", "")).strip()
        investimento_raw = str(row.get("Investimento", "")).strip().replace(".", "").replace(",", ".")

        if not canal or not investimento_raw:
            continue

        try:
            investimento = float(investimento_raw)
        except ValueError:
            continue

        if investimento < 0:
            raise ValueError("Foram encontrados valores negativos de investimento.")

        cleaned.append({"Canal": canal, "Investimento": investimento})

    if not cleaned:
        raise ValueError("Não há dados válidos após limpeza da base.")

    logging.info("Linhas válidas após limpeza: %s", len(cleaned))
    return cleaned


def build_channel_ranking(rows: Iterable[dict[str, float | str]]) -> list[RankingRow]:
    totals: dict[str, float] = defaultdict(float)
    for row in rows:
        totals[str(row["Canal"])] += float(row["Investimento"])

    ordered = sorted(totals.items(), key=lambda item: item[1], reverse=True)
    total = sum(value for _, value in ordered)
    if total <= 0:
        raise ValueError("A soma total de investimento deve ser maior que zero.")

    ranking: list[RankingRow] = []
    invest_acum = 0.0
    share_acum = 0.0

    for canal, investimento in ordered:
        share = investimento / total
        invest_acum += investimento
        share_acum += share
        ranking.append(
            RankingRow(
                canal=canal,
                investimento=investimento,
                share=share,
                share_percentual=share * 100,
                investimento_acumulado=invest_acum,
                share_acumulado=share_acum,
            )
        )
    return ranking


def compute_kpis(ranking: list[RankingRow]) -> KPIBundle:
    top = ranking[0]
    total = sum(item.investimento for item in ranking)
    media = total / len(ranking)
    hhi = sum(item.share**2 for item in ranking)
    return KPIBundle(
        total_investimento=total,
        canais_ativos=len(ranking),
        investimento_medio_por_canal=media,
        top_canal=top.canal,
        top_canal_investimento=top.investimento,
        top_canal_share=top.share,
        hhi_concentracao=hhi,
    )


def save_csv_dicts(path: Path, rows: list[dict[str, object]], headers: list[str]) -> None:
    with path.open("w", encoding="utf-8", newline="") as fp:
        writer = csv.DictWriter(fp, fieldnames=headers)
        writer.writeheader()
        writer.writerows(rows)


def save_table_outputs(clean_rows: list[dict[str, float | str]], ranking: list[RankingRow], kpis: KPIBundle, output_dir: Path) -> ReportArtifacts:
    base_processada_csv = output_dir / "base_processada.csv"
    ranking_canais_csv = output_dir / "ranking_canais.csv"
    kpis_json = output_dir / "kpis_resumo.json"
    kpis_csv = output_dir / "kpis_resumo.csv"

    save_csv_dicts(base_processada_csv, clean_rows, ["Canal", "Investimento"])
    save_csv_dicts(
        ranking_canais_csv,
        [asdict(item) for item in ranking],
        ["canal", "investimento", "share", "share_percentual", "investimento_acumulado", "share_acumulado"],
    )

    with kpis_json.open("w", encoding="utf-8") as fp:
        json.dump(asdict(kpis), fp, ensure_ascii=False, indent=2)

    save_csv_dicts(kpis_csv, [asdict(kpis)], list(asdict(kpis).keys()))

    return ReportArtifacts(
        base_processada_csv=base_processada_csv,
        ranking_canais_csv=ranking_canais_csv,
        kpis_json=kpis_json,
        kpis_csv=kpis_csv,
        grafico_pizza_png=output_dir / "investment_share_pie.png",
        grafico_barras_png=output_dir / "investment_share_bar.png",
        apresentacao_pptx=output_dir / "investment_share_report.pptx",
    )


def save_pie_chart(ranking: list[RankingRow], title: str, output_path: Path) -> None:
    fig, ax = plt.subplots(figsize=(10, 6))
    labels = [f"{item.canal} ({item.share:.1%})" for item in ranking]
    values = [item.investimento for item in ranking]
    ax.pie(values, labels=labels, autopct="%1.1f%%", startangle=90, counterclock=False)
    ax.set_title(title)
    ax.axis("equal")
    plt.tight_layout()
    fig.savefig(output_path, dpi=220)
    plt.close(fig)


def save_bar_chart(ranking: list[RankingRow], title: str, output_path: Path) -> None:
    fig, ax = plt.subplots(figsize=(11, 6))
    channels = [item.canal for item in ranking]
    values = [item.investimento for item in ranking]
    ax.bar(channels, values, color="#4472C4")
    ax.set_title(f"{title} - Investimento Absoluto")
    ax.set_xlabel("Canal")
    ax.set_ylabel("Investimento")
    ax.tick_params(axis="x", rotation=30)
    for idx, value in enumerate(values):
        ax.text(idx, value, f"{value:,.0f}", ha="center", va="bottom", fontsize=9)
    plt.tight_layout()
    fig.savefig(output_path, dpi=220)
    plt.close(fig)


def add_kpi_slide(presentation: Presentation, kpis: KPIBundle, campaign: str, period: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Resumo Executivo de KPIs"
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(12.0), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True

    lines = [
        f"Campanha: {campaign}",
        f"Período: {period}",
        f"Total investido: R$ {kpis.total_investimento:,.2f}",
        f"Canais ativos: {kpis.canais_ativos}",
        f"Investimento médio por canal: R$ {kpis.investimento_medio_por_canal:,.2f}",
        f"Canal líder: {kpis.top_canal} (R$ {kpis.top_canal_investimento:,.2f} | {kpis.top_canal_share:.1%})",
        f"HHI de concentração: {kpis.hhi_concentracao:.3f}",
    ]

    for i, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(20 if i == 0 else 16)


def add_chart_slide(presentation: Presentation, title: str, image_path: Path) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.4), width=Inches(11.7))


def save_ppt(pie_path: Path, bar_path: Path, title: str, campaign: str, period: str, kpis: KPIBundle, output_path: Path) -> None:
    presentation = Presentation()
    add_kpi_slide(presentation, kpis, campaign, period)
    add_chart_slide(presentation, title, pie_path)
    add_chart_slide(presentation, f"{title} - Visão Absoluta", bar_path)
    presentation.save(output_path)


def main() -> None:
    args = parse_args()
    setup_logging(args.log_level)

    input_path = Path(args.input)
    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    if not input_path.exists():
        raise FileNotFoundError(f"Arquivo não encontrado: {input_path}")

    clean_rows = load_and_validate_data(input_path, args.sheet)
    ranking = build_channel_ranking(clean_rows)
    kpis = compute_kpis(ranking)
    artifacts = save_table_outputs(clean_rows, ranking, kpis, output_dir)

    save_pie_chart(ranking, args.title, artifacts.grafico_pizza_png)
    save_bar_chart(ranking, args.title, artifacts.grafico_barras_png)
    save_ppt(
        artifacts.grafico_pizza_png,
        artifacts.grafico_barras_png,
        args.title,
        args.campaign,
        args.period,
        kpis,
        artifacts.apresentacao_pptx,
    )

    print("Análise concluída com sucesso.")
    print(f"Data de execução: {datetime.now().isoformat(timespec='seconds')}")
    print("Arquivos gerados:")
    for field_name, file_path in asdict(artifacts).items():
        print(f"- {field_name}: {file_path}")


if __name__ == "__main__":
    main()
